from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image


class PPTImageInserter:
    """
    A class to insert images into a PowerPoint presentation in a grid layout, maintaining the aspect ratio of images.

    Attributes
    ----------
    grid_dims : tuple
        The number of rows and columns in the grid layout.
    spacing : tuple
        The horizontal and vertical spacing between images in inches.
    title_font_size : int
        The font size of the slide title.

    Methods
    -------
    add_slide(title):
        Adds a new slide with the given title.
    add_image(image_path, title=None):
        Adds an image to the current slide, and creates a new slide if the grid is full.
    save(file_path):
        Saves the presentation to the specified file path.
    """

    def __init__(self, grid_dims=(3, 3), spacing=(0.05, 0.05), title_font_size=16):
        self.prs = Presentation()
        self.grid_rows, self.grid_cols = grid_dims
        self.spacing_horiz, self.spacing_vert = spacing
        self.title_font_size = title_font_size
        self.image_index = 0
        self.slide = None

        self.slide_width = 10
        self.slide_height = 7.5

        # Calculate title height based on font size (rough estimate)
        # Convert points to inches, with some padding
        self.title_height = Pt(self.title_font_size).pt / 72 * 1.2

        # Calculate available width and height for images
        self.available_width = self.slide_width - \
            (self.grid_cols + 1) * self.spacing_horiz
        self.available_height = self.slide_height - self.title_height - \
            0.2 - (self.grid_rows + 1) * self.spacing_vert

        # Calculate optimal image size
        self.image_width = self.available_width / self.grid_cols
        self.image_height = self.available_height / self.grid_rows

    def add_slide(self, slide_title=None):
        if slide_title is None:
            self.title_height = 0
            slide_layout = self.prs.slide_layouts[6]
            self.slide = self.prs.slides.add_slide(slide_layout)
        else:
            slide_layout = self.prs.slide_layouts[5]
            self.slide = self.prs.slides.add_slide(slide_layout)
            title_shape = self.slide.shapes.title
            title_shape.text = slide_title
            title_shape.text_frame.paragraphs[0].font.size = Pt(
                self.title_font_size)
            title_shape.left = Inches(0.5)
            title_shape.top = Inches(0.2)
            title_shape.width = Inches(8)
            title_shape.height = Inches(self.title_height)

        self.image_index = 0  # Reset image index for new slide

    def add_image(self, image_path, slide_title=None):
        if self.slide is None or self.image_index >= self.grid_rows * self.grid_cols:
            self.add_slide(slide_title)

        row = self.image_index // self.grid_cols
        col = self.image_index % self.grid_cols

        left = Inches(col * (self.image_width +
                      self.spacing_horiz) + self.spacing_horiz)
        top = Inches(self.title_height + 0.2 + row *
                     (self.image_height + self.spacing_vert) + self.spacing_vert)

        # Open the image to get its original dimensions
        with Image.open(image_path) as img:
            img_width, img_height = img.size
            img_aspect = img_width / img_height

        # Calculate new dimensions while maintaining aspect ratio
        if img_aspect > 1:
            # Wide image
            width = Inches(
                min(self.image_width, self.image_height * img_aspect))
            height = width / img_aspect
        else:
            # Tall image
            height = Inches(
                min(self.image_height, self.image_width / img_aspect))
            width = height * img_aspect

        self.slide.shapes.add_picture(
            image_path, left, top, width=width, height=height)
        self.image_index += 1

    def save(self, file_path):
        self.prs.save(file_path)


# Usage example
if __name__ == "__main__":
    import matplotlib.pyplot as plt
    import os

    # Create a directory to save images if it doesn't exist
    os.makedirs("images", exist_ok=True)

    # Generate and save some example plots
    N_images = 4
    for i in range(N_images):
        plt.figure()
        plt.plot([0, 1, 2, 3], [i, i + 1, i + 2, i + 3])
        plt.title(f"Plot {i+1}")
        plt.xlabel("X-axis")
        plt.ylabel("Y-axis")
        plt_path = f"images/plotA_{i+1}.png"
        plt.savefig(plt_path)
        plt.close()

    for i in range(N_images):
        plt.figure()
        plt.plot([0, 1, 2, 3], [i, i + 1, i + 2, i + 3], 'k')
        plt.title(f"Plot {i+1}")
        plt.xlabel("X-axis")
        plt.ylabel("Y-axis")
        plt_path = f"images/plotB_{i+1}.png"
        plt.savefig(plt_path)
        plt.close()

    # Define the paths to the images
    imageA_paths = [f"images/plotA_{i+1}.png" for i in range(N_images)]
    imageB_paths = [f"images/plotB_{i+1}.png" for i in range(N_images)]

    # Create an instance of PPTImageInserter
    ppt_inserter = PPTImageInserter(grid_dims=(2, 2))

    # Add a new slide and insert images
    for i, image_path in enumerate(imageA_paths):
        ppt_inserter.add_image(image_path, title="Hello world!")
    ppt_inserter.add_slide()
    for i, image_path in enumerate(imageB_paths):
        ppt_inserter.add_image(image_path)

    # ppt_inserter.add_slide("Plot B")
    # for i, image_path in enumerate(imageB_paths):
    #     ppt_inserter.add_image(image_path, title="Plot B")

    # Save the presentation
    ppt_inserter.save("output_presentation.pptx")
